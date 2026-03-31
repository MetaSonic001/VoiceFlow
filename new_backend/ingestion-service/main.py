from fastapi import FastAPI, HTTPException, BackgroundTasks
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import uuid

# Load .env file automatically when running locally
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass  # dotenv not installed — env vars must be set manually
import concurrent.futures
from datetime import datetime
import redis
import chromadb
from chromadb.config import Settings
import trafilatura
from bs4 import BeautifulSoup
import asyncio
from playwright.async_api import async_playwright
from pdfminer.high_level import extract_text
import boto3
from sentence_transformers import SentenceTransformer
import io
from crawl4ai import AsyncWebCrawler

# Enhanced document processing imports
from doctr.io import DocumentFile
from doctr.models import ocr_predictor
import docx
import pptx
import openpyxl
import pandas as pd
from PIL import Image
import pytesseract
import cv2
from langchain_text_splitters import RecursiveCharacterTextSplitter
from unstructured.partition.auto import partition
import tempfile
import mimetypes
import scrapy
from scrapy.crawler import CrawlerProcess
from twisted.internet import reactor, defer

app = FastAPI(title="Ingestion Service", version="1.0.0")

# Thread pool for CPU/IO-bound background tasks (keeps the event loop free)
_thread_pool = concurrent.futures.ThreadPoolExecutor(max_workers=4)

# Initialize clients
redis_client = redis.Redis(host=os.getenv("REDIS_HOST", "localhost"), port=int(os.getenv("REDIS_PORT", 6379)), decode_responses=True)
chroma_client = chromadb.HttpClient(
    host=os.getenv("CHROMA_HOST", "localhost"),
    port=int(os.getenv("CHROMA_PORT", 8002)),
)
s3_client = boto3.client(
    's3',
    endpoint_url=os.getenv("MINIO_ENDPOINT", "http://localhost:9000"),
    aws_access_key_id=os.getenv("MINIO_ACCESS_KEY"),
    aws_secret_access_key=os.getenv("MINIO_SECRET_KEY")
)

# Load embedding model
embedding_model = SentenceTransformer(os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2"))

# Initialize OCR model (doctr)
try:
    ocr_model = ocr_predictor(pretrained=True)
    print("DocTR OCR model loaded successfully")
except Exception as e:
    print(f"Failed to load DocTR OCR model: {e}")
    ocr_model = None

# Initialize text splitter for better chunking
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=int(os.getenv("CHUNK_SIZE", "1000")),
    chunk_overlap=int(os.getenv("CHUNK_OVERLAP", "200")),
    separators=["\n\n", "\n", ". ", " ", ""]
)

class IngestRequest(BaseModel):
    tenantId: str
    agentId: str
    urls: Optional[List[str]] = []
    s3_urls: Optional[List[str]] = []

class CompanyIngestRequest(BaseModel):
    tenantId: str
    website_url: str
    company_name: str
    company_description: Optional[str] = None
    industry: Optional[str] = None
    use_case: Optional[str] = None
    additional_urls: Optional[List[str]] = []
    # scraper strategy: "auto" (default), "trafilatura", "playwright", "requests"
    # auto = try trafilatura; if content < 500 chars, retry with playwright
    strategy: Optional[str] = "auto"

class IngestResponse(BaseModel):
    job_id: str
    status: str

@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "ingestion-service", "timestamp": datetime.now().isoformat()}

@app.post("/ingest", response_model=IngestResponse)
async def ingest_documents(request: IngestRequest, background_tasks: BackgroundTasks):
    job_id = str(uuid.uuid4())
    redis_client.set(f"job:{job_id}", "processing")
    redis_client.set(f"job:{job_id}:progress", "0")

    background_tasks.add_task(process_ingestion, job_id, request.tenantId, request.agentId, request.urls or [], request.s3_urls or [])

    return IngestResponse(job_id=job_id, status="processing")

@app.get("/status/{job_id}")
async def get_job_status(job_id: str):
    status = redis_client.get(f"job:{job_id}")
    progress = redis_client.get(f"job:{job_id}:progress")
    chunks_processed = redis_client.get(f"job:{job_id}:chunks")
    pages_scraped = redis_client.get(f"job:{job_id}:pages")
    return {
        "status": status,
        "progress": progress,
        "chunks_processed": int(chunks_processed) if chunks_processed else 0,
        "pages_scraped": int(pages_scraped) if pages_scraped else 0,
    }


# ─── Company-profile ingestion ────────────────────────────────────────────────
# Dedicated endpoint that scrapes a company's website and key sub-pages, then
# stores every chunk in ChromaDB tagged with source_type="company_profile".
# This keeps company knowledge clearly separated from agent-specific knowledge.
# ─── Synchronous company ingestion (runs in thread pool, never blocks the loop) ─

def _scrape_url_sync(url: str, strategy: str = "auto") -> Optional[str]:
    """
    Sync scraper with selectable strategy.
    strategy="auto"        → trafilatura first; if content < 500 chars retry with playwright
    strategy="trafilatura" → trafilatura only
    strategy="playwright"  → playwright only
    strategy="requests"    → requests+BS4 only
    """
    import requests as _req
    import urllib3
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    # If strategy is explicitly playwright, skip straight to playwright
    if strategy == "playwright":
        return _scrape_with_playwright(url)

    # Strategy 1: Trafilatura
    trafilatura_content: Optional[str] = None
    try:
        downloaded = trafilatura.fetch_url(url, no_ssl=True)
        if downloaded:
            content = trafilatura.extract(
                downloaded,
                include_comments=False,
                include_tables=True,
                include_images=False,
                include_formatting=False,
                include_links=False,
                favor_precision=True,
            )
            if content and len(content.strip()) > 100:
                trafilatura_content = content
                # For explicit trafilatura strategy, return immediately
                if strategy == "trafilatura":
                    return trafilatura_content
                # For auto: if we have enough content, no need for playwright
                if len(content.strip()) >= 500:
                    return trafilatura_content
    except Exception as e:
        print(f"Trafilatura failed for {url}: {e}")

    # If explicit requests strategy, skip to requests+BS4
    if strategy == "requests":
        pass  # fall through to requests block below

    # Strategy 2: requests + BeautifulSoup
    got_404 = False
    bs4_content: Optional[str] = None
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        r = _req.get(url, headers=headers, timeout=15, verify=False)
        if r.status_code == 404:
            got_404 = True
            return None
        r.raise_for_status()
        soup = BeautifulSoup(r.content, "html.parser")
        for el in soup(["script", "style", "nav", "header", "footer", "aside"]):
            el.decompose()
        text = soup.get_text(separator="\n", strip=True)
        lines = [l.strip() for l in text.split("\n") if l.strip() and len(l) > 10]
        result = "\n".join(lines)
        if len(result) > 100:
            bs4_content = result
            if strategy == "requests":
                return bs4_content
            # For auto: if requests found enough, no playwright needed
            if len(result) >= 500:
                return bs4_content
    except Exception as e:
        if "404" in str(e):
            got_404 = True
        else:
            print(f"requests+BS4 failed for {url}: {e}")

    # Strategy 3 / auto-fallback: Playwright (JS-heavy pages) — skip for known 404s
    # In auto mode this only runs when trafilatura AND requests returned thin content
    if not got_404 and strategy != "trafilatura" and strategy != "requests":
        playwright_result = _scrape_with_playwright(url)
        if playwright_result:
            return playwright_result

    # Return best effort from earlier strategies
    return trafilatura_content or bs4_content or None


def _scrape_with_playwright(url: str) -> Optional[str]:
    """Isolated playwright scrape — used as explicit strategy or auto-fallback."""
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            try:
                ctx = browser.new_context(ignore_https_errors=True)
                page = ctx.new_page()
                page.set_extra_http_headers({"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
                resp = page.goto(url, timeout=15000, wait_until="load")
                if resp and resp.status == 404:
                    return None
                page.wait_for_timeout(1000)
                html = page.content()
                soup = BeautifulSoup(html, "html.parser")
                for el in soup(["script", "style", "nav", "header", "footer", "aside"]):
                    el.decompose()
                text = soup.get_text(separator="\n", strip=True)
                lines = [l.strip() for l in text.split("\n") if l.strip() and len(l) > 10]
                result = "\n".join(lines)
                return result if len(result) > 100 else None
            finally:
                browser.close()
    except Exception as e:
        print(f"Playwright failed for {url}: {e}")
        return None


def _extract_internal_links(base_url: str, html: str) -> List[str]:
    """Extract up to 30 internal links from homepage HTML for deep crawling."""
    from urllib.parse import urljoin, urlparse
    try:
        base_parsed = urlparse(base_url)
        base_domain = base_parsed.netloc
        SKIP = {
            "login", "register", "signup", "cart", "checkout",
            "account", "admin", "wp-admin", "wp-login",
            ".jpg", ".jpeg", ".png", ".gif", ".svg", ".ico",
            ".css", ".js", ".woff", ".pdf", ".zip",
            "javascript:", "mailto:", "tel:", "whatsapp:",
            "facebook.com", "twitter.com", "instagram.com",
            "linkedin.com", "youtube.com", "t.me",
            "?share=", "&share=", "?print=",
        }
        soup = BeautifulSoup(html, "html.parser")
        seen = set()
        links: List[str] = []
        for a in soup.find_all("a", href=True):
            href = a["href"].strip()
            if not href or href.startswith("#"):
                continue
            if any(skip in href.lower() for skip in SKIP):
                continue
            full = urljoin(base_url, href)
            parsed = urlparse(full)
            if parsed.netloc != base_domain:
                continue
            # Normalize: strip fragment + trailing slash
            clean = parsed._replace(fragment="", query="").geturl().rstrip("/")
            if clean in seen or clean == base_url.rstrip("/"):
                continue
            seen.add(clean)
            links.append(clean)
            if len(links) >= 30:
                break
        return links
    except Exception as e:
        print(f"Link extraction error: {e}")
        return []


def _process_company_sync(
    job_id: str,
    tenant_id: str,
    website_url: str,
    company_name: str,
    additional_urls: List[str],
    company_description: Optional[str] = None,
    industry: Optional[str] = None,
    use_case: Optional[str] = None,
    strategy: str = "auto",
):
    """
    Synchronous company ingestion — runs inside a ThreadPoolExecutor so the
    FastAPI event loop is never blocked during scraping / embedding.
    """
    import requests as _req
    import urllib3
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    try:
        collection_name = f"tenant_{tenant_id}"
        collection = chroma_client.get_or_create_collection(name=collection_name)

        base = website_url.rstrip("/")

        # Step 1: Fetch homepage HTML to discover real internal links
        homepage_html = None
        try:
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            r = _req.get(base, headers=headers, timeout=15, verify=False)
            if r.ok:
                homepage_html = r.text
        except Exception as e:
            print(f"Homepage fetch failed for link discovery: {e}")

        # Step 2: Build URL list — predefined paths + discovered links from homepage
        candidate_paths = [
            "", "/about", "/about-us", "/about-us.html", "/company", "/products",
            "/services", "/solutions", "/contact", "/contact-us",
            "/team", "/our-team", "/leadership", "/careers", "/faq",
            "/index.html", "/home",
        ]
        urls_to_scrape = list(dict.fromkeys([base + p for p in candidate_paths] + list(additional_urls)))

        # Discover real internal links from homepage
        discovered: List[str] = []
        if homepage_html:
            discovered = _extract_internal_links(base, homepage_html)
            for link in discovered:
                if link not in urls_to_scrape and len(urls_to_scrape) < 30:
                    urls_to_scrape.append(link)

        # If requests-based discovery found nothing, fall back to playwright
        # (JS-heavy sites render navigation dynamically)
        if not discovered:
            try:
                from playwright.sync_api import sync_playwright
                with sync_playwright() as pw:
                    browser = pw.chromium.launch(headless=True)
                    try:
                        ctx = browser.new_context(ignore_https_errors=True)
                        page = ctx.new_page()
                        page.goto(base, timeout=15000, wait_until="load")
                        page.wait_for_timeout(1000)
                        js_html = page.content()
                        js_links = _extract_internal_links(base, js_html)
                        for link in js_links:
                            if link not in urls_to_scrape and len(urls_to_scrape) < 30:
                                urls_to_scrape.append(link)
                        print(f"Playwright discovered {len(js_links)} extra links on {base}")
                    finally:
                        browser.close()
            except Exception as e:
                print(f"Playwright link discovery failed for {base}: {e}")

        # Step 3: Optionally store structured company profile as first chunk
        if company_description or industry or use_case:
            profile_parts = [f"Company: {company_name}"]
            if industry:
                profile_parts.append(f"Industry: {industry}")
            if use_case:
                profile_parts.append(f"Primary Use Case: {use_case}")
            if website_url:
                profile_parts.append(f"Website: {website_url}")
            if company_description:
                profile_parts.append(f"Description: {company_description}")
            profile_text = "\n".join(profile_parts)
            try:
                profile_emb = embedding_model.encode([profile_text])
                collection.add(
                    ids=[f"company_profile_{tenant_id}_meta"],
                    embeddings=profile_emb.tolist(),
                    metadatas=[{
                        "agentId": "company_profile",
                        "tenantId": tenant_id,
                        "source": website_url,
                        "chunk": 0,
                        "content_type": "company_metadata",
                        "source_type": "company_profile",
                        "company_name": company_name,
                    }],
                    documents=[profile_text],
                )
                print(f"Stored company profile metadata for {company_name}")
            except Exception as e:
                print(f"Failed to store profile metadata: {e}")

        total_chunks = 0
        total_pages = 0

        for i, url in enumerate(urls_to_scrape):
            try:
                content = _scrape_url_sync(url, strategy)
                if content:
                    chunks = text_splitter.split_text(content)
                    if chunks:
                        embeddings = embedding_model.encode(chunks)
                        ids = [f"company_{tenant_id}_{uuid.uuid4()}_{j}" for j in range(len(chunks))]
                        metadatas = [
                            {
                                "agentId": "company_profile",
                                "tenantId": tenant_id,
                                "source": url,
                                "chunk": j,
                                "content_type": "company_webpage",
                                "source_type": "company_profile",
                                "company_name": company_name,
                            }
                            for j in range(len(chunks))
                        ]
                        collection.add(
                            ids=ids,
                            embeddings=embeddings.tolist(),
                            metadatas=metadatas,
                            documents=chunks,
                        )
                        total_chunks += len(chunks)
                        total_pages += 1
                        print(f"Company scrape [{company_name}] {url}: {len(chunks)} chunks")
            except Exception as e:
                print(f"Skipped {url}: {e}")

            # Always update progress — even if this URL failed/returned no content
            progress_pct = int((i + 1) / len(urls_to_scrape) * 100)
            redis_client.set(f"job:{job_id}:progress", str(progress_pct))
            redis_client.set(f"job:{job_id}:chunks", str(total_chunks))
            redis_client.set(f"job:{job_id}:pages", str(total_pages))

        redis_client.set(f"job:{job_id}", "completed")
        print(f"Company ingestion done for {company_name}: {total_pages} pages, {total_chunks} chunks")

    except Exception as e:
        redis_client.set(f"job:{job_id}", f"failed: {str(e)}")
        print(f"Company ingestion failed for {company_name}: {e}")


@app.post("/ingest/company", response_model=IngestResponse)
async def ingest_company(request: CompanyIngestRequest, background_tasks: BackgroundTasks):
    job_id = str(uuid.uuid4())
    redis_client.set(f"job:{job_id}", "processing")
    redis_client.set(f"job:{job_id}:progress", "0")
    redis_client.set(f"job:{job_id}:chunks", "0")
    redis_client.set(f"job:{job_id}:pages", "0")

    # Submit to thread pool so the event loop stays free to handle /status polls
    loop = asyncio.get_event_loop()
    loop.run_in_executor(
        _thread_pool,
        _process_company_sync,
        job_id,
        request.tenantId,
        request.website_url,
        request.company_name,
        request.additional_urls or [],
        request.company_description,
        request.industry,
        request.use_case,
        request.strategy or "auto",
    )
    return IngestResponse(job_id=job_id, status="processing")


# ─── Company knowledge browser endpoints ─────────────────────────────────────
@app.get("/knowledge/company/{tenant_id}")
async def get_company_knowledge(tenant_id: str, limit: int = 200):
    """
    Returns all ChromaDB chunks tagged as company_profile for a given tenant.
    Used by the dashboard's Knowledge Base tab to show/manage company data.
    """
    try:
        collection_name = f"tenant_{tenant_id}"
        collection = chroma_client.get_or_create_collection(name=collection_name)

        results = collection.get(
            where={"source_type": "company_profile"},
            limit=limit,
            include=["documents", "metadatas"],
        )

        chunks = []
        ids = results.get("ids") or []
        docs = results.get("documents") or []
        metas = results.get("metadatas") or []

        for chunk_id, doc, meta in zip(ids, docs, metas):
            chunks.append({"id": chunk_id, "content": doc, "metadata": meta})

        return {"chunks": chunks, "total": len(chunks)}
    except Exception as e:
        print(f"Error fetching company knowledge: {e}")
        return {"chunks": [], "total": 0}


@app.delete("/knowledge/{tenant_id}/{chunk_id}")
async def delete_knowledge_chunk(tenant_id: str, chunk_id: str):
    """Deletes a single chunk from a tenant's ChromaDB collection."""
    try:
        collection_name = f"tenant_{tenant_id}"
        collection = chroma_client.get_collection(name=collection_name)
        collection.delete(ids=[chunk_id])
        return {"deleted": True, "id": chunk_id}
    except Exception as e:
        print(f"Error deleting chunk {chunk_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

async def process_ingestion(job_id: str, tenant_id: str, agent_id: str, urls: List[str], s3_urls: List[str]):
    try:
        collection_name = f"tenant_{tenant_id}"
        collection = chroma_client.get_or_create_collection(name=collection_name)

        total_items = len(urls) + len(s3_urls)
        processed = 0

        # Process URLs
        for url in urls:
            try:
                content = await scrape_url(url)
                if content:
                    # Use semantic text splitting instead of simple chunking
                    chunks = text_splitter.split_text(content)

                    if chunks:  # Only process if we have chunks
                        embeddings = embedding_model.encode(chunks)
                        ids = [f"{agent_id}_{url}_{i}" for i in range(len(chunks))]
                        metadatas = [{"agentId": agent_id, "source": url, "chunk": i, "content_type": "webpage"} for i in range(len(chunks))]
                        collection.add(ids=ids, embeddings=embeddings.tolist(), metadatas=metadatas, documents=chunks)
                        print(f"Processed URL {url}: {len(chunks)} chunks")
            except Exception as e:
                print(f"Error processing URL {url}: {e}")

            processed += 1
            redis_client.set(f"job:{job_id}:progress", str(int(processed / total_items * 100)))

        # Process S3 URLs
        for s3_url in s3_urls:
            try:
                result = await download_from_s3(s3_url)
                if result and result['content']:
                    # Use semantic text splitting instead of simple chunking
                    chunks = text_splitter.split_text(result['content'])

                    if chunks:  # Only process if we have chunks
                        embeddings = embedding_model.encode(chunks)
                        ids = [f"{agent_id}_{s3_url}_{i}" for i in range(len(chunks))]

                        # Enhanced metadata including file info
                        metadatas = [{
                            "agentId": agent_id,
                            "source": s3_url,
                            "chunk": i,
                            "filename": result.get('filename', ''),
                            "file_type": result.get('file_type', ''),
                            **result.get('metadata', {})
                        } for i in range(len(chunks))]

                        collection.add(ids=ids, embeddings=embeddings.tolist(), metadatas=metadatas, documents=chunks)
                        print(f"Processed {result.get('filename', s3_url)}: {len(chunks)} chunks")
            except Exception as e:
                print(f"Error processing S3 URL {s3_url}: {e}")

            processed += 1
            redis_client.set(f"job:{job_id}:progress", str(int(processed / total_items * 100)))

        redis_client.set(f"job:{job_id}", "completed")

    except Exception as e:
        redis_client.set(f"job:{job_id}", f"failed: {str(e)}")

async def scrape_url(url: str) -> Optional[str]:
    """Enhanced web scraping with multiple strategies"""
    try:
        # Strategy 1: Crawl4AI (AI-driven, best for modern websites)
        try:
            async with AsyncWebCrawler() as crawler:
                result = await crawler.arun(
                    url=url,
                    config={
                        'wait_for': 'css:.content',  # Wait for content to load
                        'page_timeout': 30000,  # 30 second timeout
                        'remove_overlay_elements': True,  # Remove popups/modals
                    }
                )
                if result and result.markdown and len(result.markdown.strip()) > 200:
                    print(f"✅ Crawl4AI successfully scraped {url} ({len(result.markdown)} chars)")
                    return result.markdown
        except Exception as e:
            print(f"Crawl4AI failed for {url}: {e}")

        # Strategy 2: Trafilatura (excellent for article-style content)
        try:
            downloaded = trafilatura.fetch_url(url)
            if downloaded:
                content = trafilatura.extract(
                    downloaded,
                    include_comments=False,
                    include_tables=True,
                    include_images=False,
                    include_formatting=False,
                    include_links=False,
                    favor_precision=True
                )
                if content and len(content.strip()) > 200:
                    print(f"✅ Trafilatura successfully scraped {url} ({len(content)} chars)")
                    return content
        except Exception as e:
            print(f"Trafilatura failed for {url}: {e}")

        # Strategy 3: Playwright (for dynamic/SPA content)
        try:
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                context = await browser.new_context(
                    user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                )
                page = await context.new_page()
                await page.goto(url, wait_until="networkidle", timeout=30000)

                # Wait for content to load
                await page.wait_for_timeout(2000)

                # Try to get main content first
                main_content = await page.query_selector('main, article, .content, #content, .post')
                if main_content:
                    html = await main_content.inner_html()
                else:
                    html = await page.content()

                await browser.close()

                soup = BeautifulSoup(html, 'html.parser')

                # Remove unwanted elements
                for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside', 'advertisement']):
                    element.decompose()

                # Extract text with better formatting preservation
                text = soup.get_text(separator='\n', strip=True)

                # Clean up excessive whitespace
                lines = [line.strip() for line in text.split('\n') if line.strip()]
                text = '\n'.join(lines)

                if len(text) > 200:
                    print(f"✅ Playwright successfully scraped {url} ({len(text)} chars)")
                    return text
        except Exception as e:
            print(f"Playwright failed for {url}: {e}")

        # Strategy 4: Scrapy fallback (for complex scraping needs)
        try:
            text = await scrape_with_scrapy(url)
            if text and len(text.strip()) > 200:
                print(f"✅ Scrapy successfully scraped {url} ({len(text)} chars)")
                return text
        except Exception as e:
            print(f"Scrapy failed for {url}: {e}")

        print(f"❌ All scraping methods failed for {url}")
        return None

    except Exception as e:
        print(f"Error scraping {url}: {e}")
        return None

async def scrape_with_scrapy(url: str) -> Optional[str]:
    """Fallback web scraping using Scrapy"""
    try:
        # This is a simplified Scrapy implementation
        # In production, you'd create proper Scrapy spiders
        import requests
        from bs4 import BeautifulSoup

        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }

        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()

        soup = BeautifulSoup(response.content, 'html.parser')

        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside']):
            element.decompose()

        # Extract main content
        main_content = soup.find('main') or soup.find('article') or soup.find(class_='content') or soup.find(id='content')
        if main_content:
            text = main_content.get_text(separator='\n', strip=True)
        else:
            text = soup.get_text(separator='\n', strip=True)

        # Clean up
        lines = [line.strip() for line in text.split('\n') if line.strip() and len(line) > 10]
        return '\n'.join(lines)

    except Exception as e:
        print(f"Scrapy fallback failed: {e}")
        return None

async def download_from_s3(s3_url: str) -> Optional[Dict[str, Any]]:
    """Download and process various document types from S3 with OCR support"""
    try:
        # Parse S3 URL: s3://bucket/key
        if s3_url.startswith("s3://"):
            parts = s3_url[5:].split("/", 1)
            bucket = parts[0]
            key = parts[1]

            response = s3_client.get_object(Bucket=bucket, Key=key)
            content = response['Body'].read()
            filename = key.split('/')[-1].lower()

            # Process based on file type with better detection
            if filename.endswith('.pdf'):
                text, metadata = await process_pdf(content, filename)
            elif filename.endswith('.docx'):
                text, metadata = await process_word(content, filename)
            elif filename.endswith('.doc'):
                text, metadata = await process_word(content, filename)
            elif filename.endswith(('.pptx', '.ppt')):
                text, metadata = await process_powerpoint(content, filename)
            elif filename.endswith(('.xlsx', '.xls')):
                text, metadata = await process_excel(content, filename)
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
                text, metadata = await process_image(content, filename)
            elif filename.endswith(('.txt', '.md', '.csv', '.json', '.xml', '.html')):
                # Handle text-based files
                try:
                    text = content.decode('utf-8')
                    metadata = {
                        'file_size': len(content),
                        'encoding': 'utf-8',
                        'processing_method': 'text_decoding'
                    }
                except UnicodeDecodeError:
                    text = content.decode('latin-1', errors='ignore')
                    metadata = {
                        'file_size': len(content),
                        'encoding': 'latin-1',
                        'processing_method': 'text_decoding_fallback'
                    }
            else:
                # Try unstructured for other formats
                text, metadata = await process_unstructured(content, filename)

            return {
                'content': text,
                'metadata': metadata,
                'filename': filename,
                'file_type': filename.split('.')[-1] if '.' in filename else 'unknown'
            }
    except Exception as e:
        print(f"Error downloading from S3 {s3_url}: {e}")
        return None

async def process_pdf(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process PDF files with OCR fallback for scanned documents"""
    try:
        # Try text extraction first
        text = extract_text(io.BytesIO(content))

        # If little text extracted, likely scanned PDF - use OCR
        if len(text.strip()) < 100 and ocr_model:
            print(f"PDF {filename} appears to be scanned, using OCR...")
            text = await perform_ocr_on_pdf(content)

        metadata = {
            'page_count': len(text.split('\n\n')) if text else 0,
            'has_ocr': len(text.strip()) >= 100,
            'processing_method': 'ocr' if len(text.strip()) >= 100 else 'text_extraction'
        }

        return text, metadata
    except Exception as e:
        print(f"Error processing PDF {filename}: {e}")
        return "", {'error': str(e)}

async def process_word(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process Word documents (.docx and .doc)"""
    try:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        text_content = ""

        if filename.endswith('.docx'):
            # Handle .docx files
            doc = docx.Document(temp_file_path)

            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))

            text_content = '\n'.join(paragraphs)

            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'word_count': len(text_content.split()),
                'processing_method': 'docx_extraction',
                'format': 'docx'
            }

        else:
            # Handle .doc files using unstructured
            try:
                elements = partition(temp_file_path)
                text_parts = []
                for element in elements:
                    if hasattr(element, 'text'):
                        text_parts.append(str(element.text))
                    else:
                        text_parts.append(str(element))

                text_content = '\n'.join(text_parts)

                metadata = {
                    'elements': len(elements),
                    'processing_method': 'unstructured_doc',
                    'format': 'doc'
                }
            except Exception as e:
                print(f"Unstructured failed for .doc, using fallback: {e}")
                # Fallback: try to extract with basic text extraction
                text_content = content.decode('utf-8', errors='ignore')
                metadata = {
                    'processing_method': 'raw_text_fallback',
                    'format': 'doc',
                    'fallback': True
                }

        os.unlink(temp_file_path)

        # Clean up the text
        if text_content:
            # Remove excessive whitespace
            import re
            text_content = re.sub(r'\n\s*\n', '\n\n', text_content)
            text_content = text_content.strip()

        return text_content, metadata

    except Exception as e:
        print(f"Error processing Word document {filename}: {e}")
        return "", {'error': str(e), 'format': filename.split('.')[-1]}

async def process_powerpoint(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process PowerPoint files with detailed content extraction"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        prs = pptx.Presentation(temp_file_path)

        slide_contents = []
        slide_metadata = {}

        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_info = {
                'shapes': len(slide.shapes),
                'has_title': False,
                'has_content': False,
                'text_elements': 0
            }

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    slide_info['text_elements'] += 1

                    # Try to identify title vs content
                    if not slide_info['has_title'] and len(shape.text.strip()) < 100:
                        slide_info['has_title'] = True
                    else:
                        slide_info['has_content'] = True

            if slide_text:
                slide_content = f"\n--- Slide {i+1} ---\n" + '\n'.join(slide_text)
                slide_contents.append(slide_content)
                slide_metadata[f'slide_{i+1}'] = slide_info

        full_text = '\n'.join(slide_contents)

        # Overall metadata
        total_text_elements = sum(info['text_elements'] for info in slide_metadata.values())
        slides_with_content = sum(1 for info in slide_metadata.values() if info['has_content'])

        metadata = {
            'total_slides': len(prs.slides),
            'slides_with_content': slides_with_content,
            'total_text_elements': total_text_elements,
            'slide_details': slide_metadata,
            'word_count': len(full_text.split()),
            'processing_method': 'pptx_detailed_extraction',
            'format': 'pptx' if filename.endswith('.pptx') else 'ppt'
        }

        os.unlink(temp_file_path)
        return full_text, metadata

    except Exception as e:
        print(f"Error processing PowerPoint {filename}: {e}")
        return "", {'error': str(e)}

async def process_excel(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process Excel files with multi-sheet support"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        # Read all sheets
        excel_data = pd.read_excel(temp_file_path, sheet_name=None)

        all_text = []
        sheet_info = {}

        for sheet_name, df in excel_data.items():
            # Convert sheet to text
            sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
            sheet_text += df.to_string(index=False)
            all_text.append(sheet_text)

            # Collect sheet metadata
            sheet_info[sheet_name] = {
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': list(df.columns),
                'non_null_rows': len(df.dropna(how='all'))
            }

        full_text = '\n'.join(all_text)

        # Overall metadata
        total_rows = sum(info['rows'] for info in sheet_info.values())
        total_columns = sum(info['columns'] for info in sheet_info.values())

        metadata = {
            'sheets': list(excel_data.keys()),
            'total_rows': total_rows,
            'total_columns': total_columns,
            'sheet_details': sheet_info,
            'processing_method': 'excel_multi_sheet',
            'format': 'xlsx' if filename.endswith('.xlsx') else 'xls'
        }

        os.unlink(temp_file_path)
        return full_text, metadata

    except Exception as e:
        print(f"Error processing Excel {filename}: {e}")
        return "", {'error': str(e)}

async def process_image(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process images with OCR"""
    try:
        # Save image temporarily
        with tempfile.NamedTemporaryFile(suffix=f'.{filename.split(".")[-1]}', delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        # Use DocTR for OCR if available
        if ocr_model:
            doc = DocumentFile.from_images([temp_file_path])
            result = ocr_model(doc)
            text = result.render()
        else:
            # Fallback to Tesseract
            image = Image.open(temp_file_path)
            text = pytesseract.image_to_string(image)

        # Get image info
        img = Image.open(temp_file_path)
        metadata = {
            'width': img.width,
            'height': img.height,
            'format': img.format,
            'processing_method': 'doctr_ocr' if ocr_model else 'tesseract_ocr'
        }

        os.unlink(temp_file_path)
        return text, metadata
    except Exception as e:
        print(f"Error processing image {filename}: {e}")
        return "", {'error': str(e)}

async def process_unstructured(content: bytes, filename: str) -> tuple[str, Dict[str, Any]]:
    """Process documents using unstructured library"""
    try:
        with tempfile.NamedTemporaryFile(suffix=f'.{filename.split(".")[-1] if "." in filename else ".tmp"}', delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        elements = partition(temp_file_path)
        text = '\n'.join([str(element) for element in elements])

        metadata = {
            'elements': len(elements),
            'processing_method': 'unstructured'
        }

        os.unlink(temp_file_path)
        return text, metadata
    except Exception as e:
        print(f"Error processing with unstructured {filename}: {e}")
        # Fallback to raw text
        try:
            return content.decode('utf-8'), {'processing_method': 'raw_text', 'fallback': True}
        except:
            return "", {'error': str(e)}

async def perform_ocr_on_pdf(content: bytes) -> str:
    """Perform OCR on PDF using DocTR"""
    try:
        if not ocr_model:
            return ""

        # Save PDF temporarily
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name

        doc = DocumentFile.from_pdf(temp_file_path)
        result = ocr_model(doc)
        text = result.render()

        os.unlink(temp_file_path)
        return text
    except Exception as e:
        print(f"Error performing OCR on PDF: {e}")
        return ""

def chunk_text(text: str, source: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = ' '.join(words[start:end])
        chunks.append(chunk)
        start = end - overlap
        if start >= len(words):
            break
    return chunks

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
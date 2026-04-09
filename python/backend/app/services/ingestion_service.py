"""
Ingestion Service — Document processing pipeline.

Pipeline:
  PDF / DOCX / Images → Docling (parsing + structure) → PaddleOCR (scanned pages)
  → Post-processing (clean + chunk) → ChromaDB + BM25 index → RAG / search

Also handles URL scraping via trafilatura / httpx fallback.
"""
import asyncio
import hashlib
import io
import json
import logging
import mimetypes
import os
import re
import tempfile
import uuid
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

import chromadb
import redis
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer

logger = logging.getLogger("voiceflow.ingestion")

# ── Thread pool for CPU-bound work (keeps event loop free) ───────────────────
_thread_pool = ThreadPoolExecutor(max_workers=4)

# ── Lazy-loaded heavy models ─────────────────────────────────────────────────
_embedding_model: Optional[SentenceTransformer] = None
_docling_converter = None
_paddle_ocr = None
_chroma_client: Optional[chromadb.HttpClient] = None
_redis_client: Optional[redis.Redis] = None

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    separators=["\n\n", "\n", ". ", " ", ""],
)


def _get_embedding_model() -> SentenceTransformer:
    global _embedding_model
    if _embedding_model is None:
        model_name = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
        _embedding_model = SentenceTransformer(model_name)
        logger.info(f"[ingestion] Loaded embedding model: {model_name}")
    return _embedding_model


def _get_docling_converter():
    """Lazy-load Docling DocumentConverter."""
    global _docling_converter
    if _docling_converter is None:
        try:
            from docling.document_converter import DocumentConverter
            _docling_converter = DocumentConverter()
            logger.info("[ingestion] Docling DocumentConverter loaded")
        except Exception as e:
            logger.warning(f"[ingestion] Docling not available: {e}")
    return _docling_converter


def _get_paddle_ocr():
    """Lazy-load PaddleOCR for scanned documents."""
    global _paddle_ocr
    if _paddle_ocr is None:
        try:
            os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")
            from paddleocr import PaddleOCR
            _paddle_ocr = PaddleOCR(use_angle_cls=True, lang="en")
            logger.info("[ingestion] PaddleOCR loaded")
        except Exception as e:
            logger.warning(f"[ingestion] PaddleOCR not available: {e}")
    return _paddle_ocr


def _get_chroma() -> Optional[chromadb.HttpClient]:
    global _chroma_client
    if _chroma_client is None:
        try:
            host = os.getenv("CHROMA_HOST", "localhost")
            port = int(os.getenv("CHROMA_PORT", "8002"))
            _chroma_client = chromadb.HttpClient(host=host, port=port)
            _chroma_client.heartbeat()
            logger.info(f"[ingestion] ChromaDB connected at {host}:{port}")
        except Exception as e:
            logger.warning(f"[ingestion] ChromaDB not available: {e}")
            _chroma_client = None
    return _chroma_client


def _get_redis() -> Optional[redis.Redis]:
    global _redis_client
    if _redis_client is None:
        try:
            host = os.getenv("REDIS_HOST", "localhost")
            port = int(os.getenv("REDIS_PORT", "6379"))
            _redis_client = redis.Redis(host=host, port=port, decode_responses=True)
            _redis_client.ping()
            logger.info(f"[ingestion] Redis connected at {host}:{port}")
        except Exception as e:
            logger.warning(f"[ingestion] Redis not available: {e}")
            _redis_client = None
    return _redis_client


def _update_job_status(job_id: str, status: str, progress: int = 0, error: str = None):
    """Update job progress in Redis."""
    r = _get_redis()
    if not r:
        return
    data = {"status": status, "progress": progress, "updatedAt": datetime.now(timezone.utc).isoformat()}
    if error:
        data["error"] = error
    try:
        r.set(f"job:{job_id}", json.dumps(data), ex=3600)
    except Exception:
        pass


# ══════════════════════════════════════════════════════════════════════════════
# 1. DOCUMENT PARSING — Docling + PaddleOCR
# ══════════════════════════════════════════════════════════════════════════════

def parse_document_with_docling(file_path: str) -> str:
    """
    Parse any document (PDF, DOCX, PPTX, XLSX, images) using Docling.
    Returns extracted text. Falls back to PaddleOCR for scanned pages.
    Plain text files (.txt, .md, .csv, .json) are read directly.
    """
    ext = Path(file_path).suffix.lower()

    # Plain text files — read directly, no need for Docling
    if ext in (".txt", ".md", ".csv", ".json", ".log", ".rst"):
        try:
            return Path(file_path).read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            return ""

    converter = _get_docling_converter()
    if not converter:
        # Fallback: try basic extraction
        return _fallback_extract(file_path)

    try:
        result = converter.convert(file_path)
        text = result.document.export_to_markdown()

        # If Docling extracted very little text, might be a scanned PDF/image
        if len(text.strip()) < 50:
            ocr_text = _ocr_with_paddle(file_path)
            if ocr_text and len(ocr_text.strip()) > len(text.strip()):
                text = ocr_text

        return text.strip()
    except Exception as e:
        logger.warning(f"Docling parsing failed for {file_path}: {e}")
        # Try PaddleOCR as fallback
        ocr_text = _ocr_with_paddle(file_path)
        if ocr_text:
            return ocr_text
        return _fallback_extract(file_path)


def _ocr_with_paddle(file_path: str) -> Optional[str]:
    """Run PaddleOCR on a file (PDF pages or images)."""
    ocr = _get_paddle_ocr()
    if not ocr:
        return None

    try:
        ext = Path(file_path).suffix.lower()
        if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"):
            result = ocr.ocr(file_path, cls=True)
            lines = []
            for page in result:
                if page:
                    for line in page:
                        if line and len(line) >= 2:
                            lines.append(line[1][0])
            return "\n".join(lines)

        elif ext == ".pdf":
            # PaddleOCR handles PDFs natively
            result = ocr.ocr(file_path, cls=True)
            lines = []
            for page in result:
                if page:
                    for line in page:
                        if line and len(line) >= 2:
                            lines.append(line[1][0])
                    lines.append("\n")  # Page break
            return "\n".join(lines)

    except Exception as e:
        logger.warning(f"PaddleOCR failed for {file_path}: {e}")
    return None


def _fallback_extract(file_path: str) -> str:
    """Basic fallback extraction without Docling."""
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            from pdfminer.high_level import extract_text
            return extract_text(file_path) or ""
        elif ext in (".docx",):
            import docx
            doc = docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext in (".pptx",):
            import pptx
            prs = pptx.Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            rows = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        rows.append(" | ".join(vals))
            return "\n".join(rows)
        elif ext in (".txt", ".md", ".csv", ".json"):
            return Path(file_path).read_text(encoding="utf-8", errors="ignore")
        elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"):
            ocr_text = _ocr_with_paddle(file_path)
            return ocr_text or ""
    except Exception as e:
        logger.warning(f"Fallback extraction failed for {file_path}: {e}")
    return ""


# ══════════════════════════════════════════════════════════════════════════════
# 2. URL SCRAPING
# ══════════════════════════════════════════════════════════════════════════════

async def scrape_url(url: str) -> str:
    """Scrape text content from a URL using trafilatura with httpx fallback."""
    import httpx

    try:
        import trafilatura

        # Download page
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 VoiceFlow Ingestion Bot"})
            resp.raise_for_status()
            html = resp.text

        # Extract with trafilatura
        text = trafilatura.extract(html, include_comments=False, include_tables=True)
        if text and len(text.strip()) > 100:
            return text.strip()

        # Fallback: BeautifulSoup basic extraction
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return text.strip()

    except Exception as e:
        logger.warning(f"URL scraping failed for {url}: {e}")
        return ""


async def scrape_company_website(base_url: str, max_pages: int = 20) -> list[dict]:
    """Crawl a company website and return list of {url, content} dicts."""
    import httpx
    from urllib.parse import urljoin, urlparse

    visited = set()
    results = []
    to_visit = [base_url]
    base_domain = urlparse(base_url).netloc

    while to_visit and len(visited) < max_pages:
        url = to_visit.pop(0)
        if url in visited:
            continue
        visited.add(url)

        try:
            async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 VoiceFlow Bot"})
                resp.raise_for_status()
                html = resp.text

            import trafilatura
            text = trafilatura.extract(html, include_comments=False, include_tables=True) or ""

            if text and len(text.strip()) > 50:
                results.append({"url": url, "content": text.strip()})

            # Discover links on same domain
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, "html.parser")
            for a in soup.find_all("a", href=True):
                link = urljoin(url, a["href"])
                parsed = urlparse(link)
                if parsed.netloc == base_domain and link not in visited:
                    # Only follow html-like pages
                    path = parsed.path.lower()
                    if not any(path.endswith(ext) for ext in (".pdf", ".jpg", ".png", ".zip", ".mp4", ".mp3")):
                        to_visit.append(link.split("#")[0].split("?")[0])

        except Exception:
            continue

    return results


# ══════════════════════════════════════════════════════════════════════════════
# 3. POST-PROCESSING — Clean + Chunk
# ══════════════════════════════════════════════════════════════════════════════

def clean_text(text: str) -> str:
    """Clean extracted text: normalize whitespace, remove artifacts."""
    if not text:
        return ""
    # Normalize whitespace
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Remove common artifacts
    text = re.sub(r"Page \d+ of \d+", "", text)
    text = re.sub(r"\x0c", "\n", text)  # Form feeds
    return text.strip()


def chunk_text(text: str, source: str = "", metadata: dict = None) -> list[dict]:
    """Split text into chunks with metadata."""
    if not text or len(text.strip()) < 10:
        return []

    chunks = text_splitter.split_text(text)
    result = []
    for i, chunk in enumerate(chunks):
        chunk_meta = {
            "source": source,
            "chunk_index": i,
            "total_chunks": len(chunks),
            **(metadata or {}),
        }
        result.append({
            "content": chunk,
            "metadata": chunk_meta,
        })
    return result


# ══════════════════════════════════════════════════════════════════════════════
# 4. EMBEDDING + STORAGE — ChromaDB + BM25 Index
# ══════════════════════════════════════════════════════════════════════════════

def embed_texts(texts: list[str]) -> list[list[float]]:
    """Generate embeddings for a list of texts."""
    model = _get_embedding_model()
    embeddings = model.encode(texts, show_progress_bar=False, convert_to_numpy=True)
    return embeddings.tolist()


def store_in_chromadb(
    tenant_id: str,
    agent_id: str,
    chunks: list[dict],
    source_type: str = "document",
) -> int:
    """
    Store document chunks in ChromaDB.
    Collection: tenant_{tenantId} (tenant-isolated)
    Metadata includes agentId for per-agent filtering.
    Returns number of chunks stored.
    """
    client = _get_chroma()
    if not client:
        logger.error("ChromaDB not available, cannot store documents")
        return 0

    if not chunks:
        return 0

    collection_name = f"tenant_{tenant_id}"
    collection = client.get_or_create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"},
    )

    # Prepare data
    texts = [c["content"] for c in chunks]
    ids = []
    metadatas = []
    for i, chunk in enumerate(chunks):
        chunk_id = hashlib.sha256(
            f"{tenant_id}:{agent_id}:{chunk['metadata'].get('source', '')}:{i}:{chunk['content'][:100]}".encode()
        ).hexdigest()[:24]
        ids.append(chunk_id)
        meta = {
            "agentId": agent_id,
            "source_type": source_type,
            **{k: str(v) if not isinstance(v, (str, int, float, bool)) else v
               for k, v in chunk["metadata"].items()},
        }
        metadatas.append(meta)

    # Generate embeddings
    embeddings = embed_texts(texts)

    # Upsert in batches of 100
    batch_size = 100
    stored = 0
    for start in range(0, len(ids), batch_size):
        end = start + batch_size
        collection.upsert(
            ids=ids[start:end],
            documents=texts[start:end],
            embeddings=embeddings[start:end],
            metadatas=metadatas[start:end],
        )
        stored += len(ids[start:end])

    logger.info(f"[ingestion] Stored {stored} chunks in {collection_name} for agent {agent_id}")
    return stored


def build_bm25_index(tenant_id: str, agent_id: str) -> None:
    """
    Build a BM25 index from all documents in a tenant+agent collection.
    Stored in Redis as serialised JSON for fast retrieval at query time.
    """
    client = _get_chroma()
    r = _get_redis()
    if not client or not r:
        return

    collection_name = f"tenant_{tenant_id}"
    try:
        collection = client.get_collection(collection_name)
    except Exception:
        return

    # Fetch all documents for this agent
    try:
        results = collection.get(
            where={"agentId": agent_id},
            include=["documents", "metadatas"],
        )
    except Exception:
        # Collection might be empty or not support where filter yet
        try:
            results = collection.get(include=["documents", "metadatas"])
        except Exception:
            return

    documents = results.get("documents", [])
    ids = results.get("ids", [])
    if not documents:
        return

    # Tokenize for BM25
    tokenized = [doc.lower().split() for doc in documents]

    # Store the corpus + ids in Redis so we can reconstruct BM25 at query time
    bm25_data = {
        "ids": ids,
        "documents": documents,
        "metadatas": results.get("metadatas", []),
        "tokenized": tokenized,
    }
    key = f"bm25:{tenant_id}:{agent_id}"
    r.set(key, json.dumps(bm25_data), ex=86400)  # 24h TTL
    logger.info(f"[bm25] Built index for {tenant_id}/{agent_id}: {len(documents)} docs")


# ══════════════════════════════════════════════════════════════════════════════
# 5. MAIN INGESTION PIPELINES
# ══════════════════════════════════════════════════════════════════════════════

async def ingest_file(
    file_bytes: bytes,
    filename: str,
    tenant_id: str,
    agent_id: str,
    job_id: str = None,
) -> dict:
    """
    Full file ingestion pipeline:
    1. Save to temp file
    2. Parse with Docling (+ PaddleOCR for scans)
    3. Clean + chunk
    4. Embed + store in ChromaDB
    5. Build BM25 index
    """
    if job_id:
        _update_job_status(job_id, "processing", 10)

    # Save to temp file
    ext = Path(filename).suffix.lower() or ".pdf"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        # Parse document (CPU-bound, run in thread pool)
        loop = asyncio.get_event_loop()
        if job_id:
            _update_job_status(job_id, "processing", 30, None)

        text = await loop.run_in_executor(
            _thread_pool, parse_document_with_docling, tmp_path
        )

        if not text or len(text.strip()) < 10:
            return {"status": "empty", "chunks": 0, "message": "No text extracted from document"}

        if job_id:
            _update_job_status(job_id, "processing", 50)

        # Clean
        text = clean_text(text)

        # Chunk
        chunks = chunk_text(text, source=filename, metadata={
            "filename": filename,
            "content_type": mimetypes.guess_type(filename)[0] or "application/octet-stream",
        })

        if job_id:
            _update_job_status(job_id, "processing", 70)

        # Store in ChromaDB (CPU-bound embedding)
        stored = await loop.run_in_executor(
            _thread_pool,
            store_in_chromadb,
            tenant_id, agent_id, chunks, "file_upload",
        )

        # Build BM25 index
        await loop.run_in_executor(
            _thread_pool,
            build_bm25_index,
            tenant_id, agent_id,
        )

        if job_id:
            _update_job_status(job_id, "completed", 100)

        return {
            "status": "completed",
            "chunks": stored,
            "characters": len(text),
            "filename": filename,
        }

    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


async def ingest_urls(
    urls: list[str],
    tenant_id: str,
    agent_id: str,
    job_id: str = None,
) -> dict:
    """
    Ingest content from URLs:
    1. Scrape each URL
    2. Clean + chunk
    3. Embed + store in ChromaDB
    4. Build BM25 index
    """
    if job_id:
        _update_job_status(job_id, "processing", 10)

    total_chunks = 0
    processed = 0
    errors = []

    for i, url in enumerate(urls):
        try:
            text = await scrape_url(url)
            if not text or len(text.strip()) < 50:
                errors.append(f"{url}: no content extracted")
                continue

            text = clean_text(text)
            chunks = chunk_text(text, source=url, metadata={"url": url})

            loop = asyncio.get_event_loop()
            stored = await loop.run_in_executor(
                _thread_pool,
                store_in_chromadb,
                tenant_id, agent_id, chunks, "url_scrape",
            )
            total_chunks += stored
            processed += 1

        except Exception as e:
            errors.append(f"{url}: {str(e)}")

        if job_id:
            progress = 10 + int(80 * (i + 1) / len(urls))
            _update_job_status(job_id, "processing", progress)

    # Build BM25 index after all URLs processed
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(
        _thread_pool,
        build_bm25_index,
        tenant_id, agent_id,
    )

    if job_id:
        _update_job_status(job_id, "completed", 100)

    return {
        "status": "completed",
        "urls_processed": processed,
        "total_chunks": total_chunks,
        "errors": errors if errors else None,
    }


async def ingest_company_website(
    website_url: str,
    tenant_id: str,
    agent_id: str,
    job_id: str = None,
    max_pages: int = 20,
) -> dict:
    """
    Crawl a company website and ingest all pages.
    """
    if job_id:
        _update_job_status(job_id, "processing", 5)

    pages = await scrape_company_website(website_url, max_pages=max_pages)

    if not pages:
        if job_id:
            _update_job_status(job_id, "completed", 100)
        return {"status": "completed", "pages": 0, "chunks": 0}

    total_chunks = 0
    for i, page in enumerate(pages):
        text = clean_text(page["content"])
        chunks = chunk_text(text, source=page["url"], metadata={"url": page["url"]})

        if chunks:
            loop = asyncio.get_event_loop()
            stored = await loop.run_in_executor(
                _thread_pool,
                store_in_chromadb,
                tenant_id, agent_id, chunks, "company_website",
            )
            total_chunks += stored

        if job_id:
            progress = 5 + int(90 * (i + 1) / len(pages))
            _update_job_status(job_id, "processing", progress)

    # Build BM25 index
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(
        _thread_pool,
        build_bm25_index,
        tenant_id, agent_id,
    )

    if job_id:
        _update_job_status(job_id, "completed", 100)

    return {
        "status": "completed",
        "pages": len(pages),
        "chunks": total_chunks,
    }


async def ingest_s3_file(
    s3_path: str,
    tenant_id: str,
    agent_id: str,
    job_id: str = None,
) -> dict:
    """Download file from MinIO/S3 and run the file ingestion pipeline."""
    import boto3

    try:
        endpoint = os.getenv("MINIO_ENDPOINT", "localhost:9000")
        if not endpoint.startswith("http"):
            endpoint = f"http://{endpoint}"

        s3 = boto3.client(
            "s3",
            endpoint_url=endpoint,
            aws_access_key_id=os.getenv("MINIO_ACCESS_KEY", "minioadmin"),
            aws_secret_access_key=os.getenv("MINIO_SECRET_KEY", "minioadmin"),
        )

        # Parse bucket and key from path: bucket/key or just key
        parts = s3_path.lstrip("/").split("/", 1)
        bucket = parts[0] if len(parts) > 1 else "voiceflow-documents"
        key = parts[1] if len(parts) > 1 else parts[0]

        obj = s3.get_object(Bucket=bucket, Key=key)
        file_bytes = obj["Body"].read()
        filename = Path(key).name

        return await ingest_file(file_bytes, filename, tenant_id, agent_id, job_id)

    except Exception as e:
        logger.exception(f"S3 ingestion failed for {s3_path}")
        if job_id:
            _update_job_status(job_id, "failed", 0, str(e))
        return {"status": "failed", "error": str(e)}


def get_job_status(job_id: str) -> dict:
    """Get job status from Redis."""
    r = _get_redis()
    if not r:
        return {"status": "unknown", "progress": 0}
    try:
        data = r.get(f"job:{job_id}")
        if data:
            return json.loads(data)
    except Exception:
        pass
    return {"status": "unknown", "progress": 0}


def delete_agent_documents(tenant_id: str, agent_id: str) -> int:
    """Delete all documents for an agent from ChromaDB and BM25 index."""
    client = _get_chroma()
    if not client:
        return 0

    collection_name = f"tenant_{tenant_id}"
    try:
        collection = client.get_collection(collection_name)
        # Get IDs for this agent
        results = collection.get(where={"agentId": agent_id}, include=[])
        ids = results.get("ids", [])
        if ids:
            collection.delete(ids=ids)
            logger.info(f"[ingestion] Deleted {len(ids)} chunks for agent {agent_id}")

        # Clear BM25 index
        r = _get_redis()
        if r:
            r.delete(f"bm25:{tenant_id}:{agent_id}")

        return len(ids)
    except Exception as e:
        logger.warning(f"Failed to delete agent documents: {e}")
        return 0
